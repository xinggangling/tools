import os
import json
import requests
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from tqdm import tqdm

class OllamaTranslator:
    def __init__(self, model_name="llama3.2:latest", base_url="http://localhost:11434"):
        self.model_name = model_name
        self.base_url = base_url
        self.api_url = f"{base_url}/api/generate"
    
    def protect_special_chars(self, text):
        """保护特殊字符，将其替换为占位符"""
        # 保护代码块
        code_blocks = []
        def replace_code(match):
            code_blocks.append(match.group(0))
            return f"__CODE_BLOCK_{len(code_blocks)-1}__"
        
        # 保护代码块（包括缩进）
        text = re.sub(r'(?m)^\s*(?:if|for|while|def|class|try|except|else|elif|with|import|from|return|print|#).*$', replace_code, text)
        
        # 保护特殊符号
        special_chars = []
        def replace_special(match):
            special_chars.append(match.group(0))
            return f"__SPECIAL_CHAR_{len(special_chars)-1}__"
        
        # 保护常见的特殊符号
        text = re.sub(r'[<>{}[\]()@#$%^&*+=|\\/~`]', replace_special, text)
        
        return text, code_blocks, special_chars
    
    def restore_special_chars(self, text, code_blocks, special_chars):
        """恢复特殊字符"""
        # 恢复代码块
        for i, code in enumerate(code_blocks):
            text = text.replace(f"__CODE_BLOCK_{i}__", code)
        
        # 恢复特殊符号
        for i, char in enumerate(special_chars):
            text = text.replace(f"__SPECIAL_CHAR_{i}__", char)
        
        return text
    
    def translate(self, text, target_lang='zh-cn'):
        """使用Ollama API翻译文本"""
        if not text or not text.strip():
            return text
        
        try:
            # 保护特殊字符
            protected_text, code_blocks, special_chars = self.protect_special_chars(text)
            
            # 构建提示词，确保只返回翻译结果
            prompt = f"""Translate the following text to {target_lang}. 
IMPORTANT: Return ONLY the translation, no explanations, no thinking process, no additional text.

Rules:
1. Keep code snippets unchanged
2. Maintain original formatting and indentation
3. Preserve technical terms accurately
4. Return ONLY the translated text
5. Keep the translation concise and similar in length to the original
6. Do not translate any text between __CODE_BLOCK_X__ and __SPECIAL_CHAR_X__ markers

Text to translate:
{protected_text}"""
            
            # 调用Ollama API
            response = requests.post(
                self.api_url,
                json={
                    "model": self.model_name,
                    "prompt": prompt,
                    "stream": False,
                    "temperature": 0.1  # 降低随机性，使输出更稳定
                }
            )
            
            if response.status_code == 200:
                result = response.json()
                translated_text = result.get('response', protected_text).strip()
                # 清理可能的额外文本
                if "Translation:" in translated_text:
                    translated_text = translated_text.split("Translation:")[-1].strip()
                if "Here's the translation:" in translated_text:
                    translated_text = translated_text.split("Here's the translation:")[-1].strip()
                
                # 恢复特殊字符
                translated_text = self.restore_special_chars(translated_text, code_blocks, special_chars)
                return translated_text
            else:
                print(f"API调用失败: {response.status_code}")
                return text
                
        except Exception as e:
            print(f"翻译出错: {e}")
            return text

def adjust_text_size(shape, translated_text):
    """调整文本大小以适应形状"""
    if not hasattr(shape, "text_frame"):
        return
    
    text_frame = shape.text_frame
    if not text_frame.text:
        return
    
    # 获取原始文本大小
    original_size = None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                original_size = run.font.size
                break
        if original_size:
            break
    
    if not original_size:
        return
    
    # 计算文本长度比例
    original_length = len(text_frame.text)
    translated_length = len(translated_text)
    length_ratio = translated_length / original_length if original_length > 0 else 1
    
    # 根据长度比例调整字体大小
    if length_ratio > 1.2:  # 如果翻译后的文本长度超过原文本20%
        new_size = int(original_size.pt / length_ratio)
        # 确保字体大小不小于8pt
        new_size = max(8, new_size)
        
        # 应用新的字体大小
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(new_size)

def translate_ppt(input_file, output_file, target_lang='zh-cn', model_name="llama3.2:latest"):
    """翻译PPT文件"""
    # 初始化翻译器
    translator = OllamaTranslator(model_name=model_name)
    
    # 加载PPT文件
    prs = Presentation(input_file)
    
    # 遍历所有幻灯片
    for slide in tqdm(prs.slides, desc="正在翻译幻灯片"):
        # 遍历所有形状
        for shape in slide.shapes:
            # 处理文本框
            if hasattr(shape, "text"):
                original_text = shape.text
                translated_text = translator.translate(original_text, target_lang)
                shape.text = translated_text
                adjust_text_size(shape, translated_text)
            
            # 处理表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            original_text = cell.text
                            translated_text = translator.translate(original_text, target_lang)
                            cell.text = translated_text
                            # 为表格单元格创建临时形状来调整文本大小
                            temp_shape = cell._tc.get_or_add_txBody()
                            adjust_text_size(temp_shape, translated_text)
    
    # 保存翻译后的文件
    prs.save(output_file)
    print(f"翻译完成！文件已保存为: {output_file}")

def main():
    # 获取用户输入
    input_file = input("请输入要翻译的PPT文件路径: ")
    output_file = input("请输入保存翻译后文件的路径: ")
    target_lang = input("请输入目标语言(默认中文): ") or '中文'
    model_name = input("请输入要使用的模型名称(默认llama3.2:latest): ") or 'llama3.2:latest'
    
    # 检查输入文件是否存在
    if not os.path.exists(input_file):
        print("错误：输入文件不存在！")
        return
    
    # 执行翻译
    translate_ppt(input_file, output_file, target_lang, model_name)

if __name__ == "__main__":
    main() 