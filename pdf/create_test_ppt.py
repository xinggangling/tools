from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_test_ppt(output_file):
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 添加标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Python编程技术文档"
    subtitle.text = "测试文档 - 2024年"
    
    # 添加目录页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "目录"
    content.text = "1. Python基础语法\n2. 面向对象编程\n3. 文件操作\n4. 异常处理\n5. 模块和包"
    
    # 添加Python基础语法页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Python基础语法"
    content.text = """变量和数据类型：
- 整数(int)：age = 25
- 浮点数(float)：height = 1.75
- 字符串(str)：name = "Python"
- 布尔值(bool)：is_student = True

条件语句：
if age >= 18:
    print("成年人")
else:
    print("未成年人")"""
    
    # 添加面向对象编程页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "面向对象编程"
    content.text = """类的定义：
class Student:
    def __init__(self, name, age):
        self.name = name
        self.age = age
    
    def study(self):
        print(f"{self.name}正在学习")

继承：
class GraduateStudent(Student):
    def research(self):
        print(f"{self.name}正在做研究")"""
    
    # 添加文件操作页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "文件操作"
    content.text = """文件读写：
# 写入文件
with open('test.txt', 'w') as f:
    f.write('Hello, World!')

# 读取文件
with open('test.txt', 'r') as f:
    content = f.read()
    print(content)"""
    
    # 添加异常处理页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "异常处理"
    content.text = """try-except语句：
try:
    number = int(input("请输入一个数字："))
    result = 100 / number
except ValueError:
    print("输入的不是有效数字")
except ZeroDivisionError:
    print("不能除以零")
else:
    print(f"结果是：{result}")"""
    
    # 添加模块和包页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "模块和包"
    content.text = """导入模块：
import math
from datetime import datetime

使用模块：
print(math.pi)  # 3.141592653589793
print(datetime.now())  # 当前时间

创建包：
my_package/
    __init__.py
    module1.py
    module2.py"""
    
    # 保存文件
    prs.save(output_file)
    print(f"测试PPT已创建：{output_file}")

if __name__ == "__main__":
    create_test_ppt("test_ppt.pptx") 